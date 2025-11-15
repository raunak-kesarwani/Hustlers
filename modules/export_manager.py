"""
Export Manager Module
Handles exporting content to various formats: PDF, PPT, DOCX
"""

import os
from datetime import datetime
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
import json

class ExportManager:
    """Manages exporting content to different file formats"""
    
    def __init__(self):
        """Initialize export manager"""
        # Create exports directory if it doesn't exist
        self.exports_dir = 'exports'
        if not os.path.exists(self.exports_dir):
            os.makedirs(self.exports_dir)
    
    def export(self, content, format_type='pdf', filename='export'):
        """
        Export content to specified format
        format_type: 'pdf', 'ppt', 'docx'
        """
        if format_type == 'pdf':
            return self._export_pdf(content, filename)
        elif format_type == 'ppt':
            return self._export_ppt(content, filename)
        elif format_type == 'docx':
            return self._export_docx(content, filename)
        else:
            raise ValueError(f"Unsupported format: {format_type}")
    
    def _export_pdf(self, content, filename):
        """Export content to PDF"""
        filepath = os.path.join(self.exports_dir, f"{filename}.pdf")
        
        # Create PDF document
        doc = SimpleDocTemplate(filepath, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        # Extract content based on type
        if isinstance(content, dict):
            # Add title
            title = content.get('topic', 'Educational Content')
            story.append(Paragraph(title, styles['Title']))
            story.append(Spacer(1, 12))
            
            # Add content
            content_text = content.get('content', content.get('script', content.get('description', str(content))))
            # Clean HTML tags if present
            content_text = content_text.replace('<p>', '').replace('</p>', '\n')
            content_text = content_text.replace('<br>', '\n')
            
            # Split into paragraphs
            paragraphs = content_text.split('\n\n')
            for para in paragraphs:
                if para.strip():
                    story.append(Paragraph(para.strip(), styles['Normal']))
                    story.append(Spacer(1, 6))
        else:
            # Simple text content
            story.append(Paragraph(str(content), styles['Normal']))
        
        # Build PDF
        doc.build(story)
        return filepath
    
    def _export_ppt(self, content, filename):
        """Export content to PowerPoint"""
        filepath = os.path.join(self.exports_dir, f"{filename}.pptx")
        
        # Create presentation
        prs = Presentation()
        
        # Extract content
        if isinstance(content, dict):
            title = content.get('topic', 'Educational Content')
            content_text = content.get('content', content.get('script', str(content)))
        else:
            title = 'Educational Content'
            content_text = str(content)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title_shape.text = title
        subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}"
        
        # Content slides
        # Split content into chunks for multiple slides
        chunks = content_text.split('\n\n')
        for chunk in chunks:
            if chunk.strip():
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes
                
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                
                title_shape.text = "Content"
                tf = body_shape.text_frame
                tf.text = chunk.strip()
        
        prs.save(filepath)
        return filepath
    
    def _export_docx(self, content, filename):
        """Export content to Word document"""
        filepath = os.path.join(self.exports_dir, f"{filename}.docx")
        
        # Create document
        doc = Document()
        
        # Extract content
        if isinstance(content, dict):
            title = content.get('topic', 'Educational Content')
            content_text = content.get('content', content.get('script', str(content)))
        else:
            title = 'Educational Content'
            content_text = str(content)
        
        # Add title
        doc.add_heading(title, 0)
        
        # Add content
        # Split into paragraphs
        paragraphs = content_text.split('\n\n')
        for para in paragraphs:
            if para.strip():
                doc.add_paragraph(para.strip())
        
        # Save document
        doc.save(filepath)
        return filepath

